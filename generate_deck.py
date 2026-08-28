import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import shutil

# Make a fresh copy from backup
if os.path.exists('CIT Hackthon_Original.pptx'):
    shutil.copyfile('CIT Hackthon_Original.pptx', 'CIT Hackthon.pptx')
else:
    shutil.copyfile('CIT Hackthon.pptx', 'CIT Hackthon_Original.pptx')

prs = pptx.Presentation('CIT Hackthon.pptx')

# Color palette
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(56, 189, 248)        # #38BDF8 (Primary accent)
LIGHT_BLUE = RGBColor(147, 197, 253) # #93C5FD
SLATE_GRAY = RGBColor(226, 232, 240) # #E2E8F0 (Subtext/body)
GREEN = RGBColor(52, 211, 153)       # #34D399 (Success/Accent)
YELLOW = RGBColor(251, 191, 36)      # #FBBF24 (Highlight)
PINK = RGBColor(244, 114, 182)       # #F472B6 (Alert)

def clear_and_set_title(shape, title_text, title_size=42):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        r = p.runs[0]
        r.font.name = "Segoe UI"
        r.font.size = Pt(title_size)
        r.font.bold = True
        r.font.color.rgb = WHITE

def set_subtitle(shape, text, font_size=20, color=CYAN):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        r = p.runs[0]
        r.font.name = "Segoe UI"
        r.font.size = Pt(font_size)
        r.font.bold = False
        r.font.italic = True
        r.font.color.rgb = color

def add_bullet(tf, heading, body, head_color=CYAN, body_color=SLATE_GRAY, head_size=18, body_size=16.5, space_after=12):
    p = tf.add_paragraph() if tf.paragraphs and tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(space_after)
    p.line_spacing = 1.18
    p.alignment = PP_ALIGN.LEFT
    
    # Heading
    r1 = p.add_run()
    r1.text = heading + (" — " if heading else "")
    r1.font.name = "Segoe UI"
    r1.font.size = Pt(head_size)
    r1.font.bold = True
    r1.font.color.rgb = head_color
    
    # Body
    r2 = p.add_run()
    r2.text = body
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(body_size)
    r2.font.bold = False
    r2.font.color.rgb = body_color

# ==============================================================================
# SLIDE 1: Title Slide
# ==============================================================================
slide1 = prs.slides[0]
for s in slide1.shapes:
    if s.name == 'object 3':
        tf = s.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = "TradeSentinel"
        p1.alignment = PP_ALIGN.LEFT
        if p1.runs:
            r = p1.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(54)
            r.font.bold = True
            r.font.color.rgb = WHITE
            
        p2 = tf.add_paragraph()
        p2.text = "AI-Powered Business Automation Copilot for Logistics"
        p2.space_before = Pt(8)
        if p2.runs:
            r = p2.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = CYAN
            
        p3 = tf.add_paragraph()
        p3.text = "Transforming Complex Supply Chain Decisions into Executable, Self-Optimizing Workflows (PS4)"
        p3.space_before = Pt(6)
        if p3.runs:
            r = p3.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(17)
            r.font.italic = True
            r.font.color.rgb = SLATE_GRAY
            
    elif s.name == 'object 2':
        tf = s.text_frame
        tf.clear()
        
        lines = [
            ("Team Name", ":  MARK42", CYAN, True),
            ("Team Leader Name", ":  GOUTHAM REDDY ESAMBADI", WHITE, True),
            ("Team Members", ":  KANIKSHA S", WHITE, True),
            ("", "   ELISETTY SOWMYA", WHITE, False),
            ("", "   DUTALLURI KARIMULLA", WHITE, False),
        ]
        
        for i, (label, val, col, is_bold) in enumerate(lines):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.space_after = Pt(6)
            if label:
                r1 = p.add_run()
                r1.text = f"{label:<18}"
                r1.font.name = "Segoe UI"
                r1.font.size = Pt(22)
                r1.font.bold = True
                r1.font.color.rgb = CYAN
            else:
                r1 = p.add_run()
                r1.text = f"{' ':<18}"
                r1.font.name = "Segoe UI"
                r1.font.size = Pt(22)
                
            r2 = p.add_run()
            r2.text = val
            r2.font.name = "Segoe UI"
            r2.font.size = Pt(22)
            r2.font.bold = is_bold
            r2.font.color.rgb = col

# ==============================================================================
# SLIDE 2: Problem Statement
# ==============================================================================
slide2 = prs.slides[1]
for s in slide2.shapes:
    if s.name == 'object 3':
        clear_and_set_title(s, "Problem Statement", title_size=42)
    elif s.name == 'object 6':
        set_subtitle(s, "The Critical Operational & Automation Bottlenecks in Modern Logistics (PS4)")
    elif s.name == 'TextBox 6':
        tf = s.text_frame
        tf.clear()
        
        add_bullet(tf, "Repetitive Manual Decision Loops", 
                   "Logistics teams repeatedly evaluate risk levels, recalculate ETAs, inspect customs delays, and coordinate rerouting across disconnected spreadsheets and portals over 180+ times per month.",
                   head_color=CYAN)
        
        add_bullet(tf, "High Technical Barrier to Automation", 
                   "Operations managers clearly understand business rules (e.g., 'If risk > 70 and delay > 2 days, find alternate routes'), but lack programming, RPA scripting, or API expertise to build automations.",
                   head_color=CYAN)
        
        add_bullet(tf, "Fragile Execution & Lack of Governance", 
                   "Traditional RPA scripts break on minor data/UI changes, lack semantic understanding, and cannot detect multi-rule collisions or enforce financial approval gates (risking unauthorized rerouting on high-value cargo > ₹10L).",
                   head_color=CYAN)
        
        add_bullet(tf, "Opaque Black-Boxes & No Continuous Optimization", 
                   "Existing automation tools cannot simulate workflows against historical operations data, lack factor-level decision explainability ('Why was it rerouted?'), and never learn from operator approvals to optimize bottlenecks.",
                   head_color=CYAN)

# ==============================================================================
# SLIDE 3: Proposed Solution
# ==============================================================================
slide3 = prs.slides[2]
for s in slide3.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "Proposed Solution: TradeSentinel Copilot", title_size=42)
    elif s.name == 'object 3':
        set_subtitle(s, "Turning Logistics Intelligence into Autonomous, Executable Workflows with Human-in-the-Loop")
    elif s.name == 'TextBox 6':
        tf = s.text_frame
        tf.clear()
        
        add_bullet(tf, "AI-Native Business Automation Copilot", 
                   "An intelligent enterprise platform that accepts workflow requirements in natural language and converts them into structured, validated, and executable workflow Directed Acyclic Graphs (DAGs).",
                   head_color=GREEN)
        
        add_bullet(tf, "Logistics Intelligence as AI Action Tools", 
                   "Transforms TradeSentinel's deterministic ML models (Risk Scoring, ETA Predictor, Customs Analyzer, Route Optimizer, Disruption Simulator) into callable tools for autonomous agentic reasoning.",
                   head_color=GREEN)
        
        add_bullet(tf, "Human-in-the-Loop Governance & Tiered Approval", 
                   "Enables autonomous straight-through execution for routine, low-risk logistics actions while enforcing mandatory human authorization gates for high-value cargo (> ₹10L) and critical route changes.",
                   head_color=GREEN)
        
        add_bullet(tf, "Closed-Loop Safety & Continuous Improvement", 
                   "Unites natural language synthesis, deterministic conflict detection, historical dataset simulation, proactive opportunity mining, and self-optimizing workflow intelligence into one coherent platform.",
                   head_color=GREEN)

# ==============================================================================
# SLIDE 4: Key Features
# ==============================================================================
slide4 = prs.slides[3]
for s in slide4.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "Key Features of TradeSentinel Copilot", title_size=42)
    elif s.name == 'object 3':
        set_subtitle(s, "6 Core Capabilities Bridging Natural Language Intent to Autonomous Execution")
    elif s.name == 'TextBox 6':
        tf = s.text_frame
        tf.clear()
        
        features = [
            ("1. Natural Language Workflow Generator", "Conversational AI parses triggers, entities, conditions, thresholds, branches, and actions into structured Pydantic workflow schemas with instant visual rendering."),
            ("2. Deterministic Conflict Detection Engine", "Identifies approval bypasses, trigger collisions, contradictory actions, race conditions, and infinite loops before deployment."),
            ("3. Historical Workflow Simulator", "Replays workflows across 500+ historical shipments to quantify delay reduction (2.4 days saved) and cost impact (₹4.8L savings) prior to live execution."),
            ("4. Proactive Automation Opportunity Detector", "Mines operational logs to discover repetitive manual procedures (e.g., 184 manual escalations/month) and suggests 1-click automated workflows."),
            ("5. Self-Optimizing Analytics & Bottleneck Engine", "Tracks execution latency and operator approval patterns to recommend rule optimizations (e.g., auto-approving reroutes < ₹5L based on 92% approval rate)."),
            ("6. Transparent AI Explainability ('Explain WHY')", "Provides factor-by-factor risk contribution breakdowns (+18 Port, +16 Geopolitical) and policy audit trails for every automated action taken.")
        ]
        
        for title, desc in features:
            add_bullet(tf, title, desc, head_color=CYAN, body_size=15, head_size=16, space_after=8)

# ==============================================================================
# SLIDE 5: Idea / Approach & Technology Stack
# ==============================================================================
slide5 = prs.slides[4]
for s in slide5.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "Idea / Approach & Technology Stack", title_size=42)
    elif s.name == 'object 3':
        set_subtitle(s, "Enterprise-Grade Multi-Tiered Architecture with Domain Logistics Intelligence")
    elif s.name == 'TextBox 6':
        tf = s.text_frame
        tf.clear()
        
        add_bullet(tf, "AI & Agentic Reasoning Layer", 
                   "LLM Workflow Parser (Claude Sonnet / Gemini / Local Ollama) enforcing strict Pydantic DSL schemas; Function calling integration binding TradeSentinel's deterministic ML tool registry.",
                   head_color=LIGHT_BLUE, body_size=15.5, head_size=16.5, space_after=10)
        
        add_bullet(tf, "Workflow DSL & Graph Execution Engine", 
                   "Directed Acyclic Graph (DAG) state machine supporting Triggers, Multi-variable Conditions, Parallel Branches, Time Delays, Manager Approvals, and Step-by-Step State Rollback Logging.",
                   head_color=LIGHT_BLUE, body_size=15.5, head_size=16.5, space_after=10)
        
        add_bullet(tf, "Deterministic Conflict & Graph Analysis Engine", 
                   "Graph cycle detection (Tarjan's algorithm), logical condition satisfiability (Z3 / boolean rule solvers), and state-transition constraint checks to ensure 100% runtime safety.",
                   head_color=LIGHT_BLUE, body_size=15.5, head_size=16.5, space_after=10)
        
        add_bullet(tf, "Full-Stack Enterprise Implementation", 
                   "Backend: Python FastAPI, Embedded Persistent VectorDB (Dense Embeddings & Cosine Search), Pydantic v2, PyJWT RBAC.\nFrontend: React 19, Vite, Tailwind CSS, React Flow (interactive node canvas), Recharts, React-Leaflet GIS.",
                   head_color=LIGHT_BLUE, body_size=15.5, head_size=16.5, space_after=10)

# ==============================================================================
# SLIDE 6: System Architecture / Workflow Diagram
# ==============================================================================
slide6 = prs.slides[5]
for s in slide6.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "System Architecture & End-to-End Workflow", title_size=42)

# Replace Picture 9 with tradesentinel_architecture.png
for s in list(slide6.shapes):
    if s.name == 'Picture 9' or (s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.width > Inches(8)):
        left, top, width, height = s.left, s.top, s.width, s.height
        sp = s._element
        sp.getparent().remove(sp)
        slide6.shapes.add_picture('tradesentinel_architecture.png', left, top, width, height)
        break

# ==============================================================================
# SLIDE 7: Innovation & Existing Solutions
# ==============================================================================
slide7 = prs.slides[6]
for s in slide7.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "Innovation & Key Differentiators", title_size=42)
    elif s.name == 'object 3':
        set_subtitle(s, "How TradeSentinel Outperforms Traditional Automation and Generic AI Chatbots")
    elif s.name == 'TextBox 6':
        tf = s.text_frame
        tf.clear()
        
        innovations = [
            ("1. Conversational NL to Executable Workflow Graph", "Unlike brittle RPA tools that record mouse clicks, TradeSentinel semantically extracts triggers, multi-variable conditions, thresholds, and actions into structured, validated Pydantic DAGs."),
            ("2. Logistics ML Intelligence as Agentic Tools", "Instead of generic AI hallucinations, our copilot directly executes domain-specific ML algorithms: risk scoring, ETA prediction, customs bottlenecks, and dynamic route optimization."),
            ("3. Deterministic Pre-Deployment Conflict Engine", "Prevents catastrophic runtime collisions, approval bypasses (e.g., auto-reroute bypassing ₹10L manager policy), contradictory actions, and infinite loops using graph/rule validation."),
            ("4. Historical Workflow Simulation & Continuous Self-Optimization", "Validates prospective workflows against 500+ historical shipments with projected ROI, and continuously analyzes operator approvals to auto-tune approval thresholds.")
        ]
        
        for title, desc in innovations:
            add_bullet(tf, title, desc, head_color=GREEN, body_size=15, head_size=16, space_after=9)

# ==============================================================================
# SLIDE 8: Impact & Future Scope (Table & Metrics)
# ==============================================================================
slide8 = prs.slides[7]
for s in slide8.shapes:
    if s.name == 'object 2':
        clear_and_set_title(s, "Measurable Impact & Future Scope", title_size=42)
        
    elif s.has_table:
        table = s.table
        
        # 4-column comparison table (7 rows: 1 header + 6 data rows)
        headers = ["Capability / Feature", "Traditional RPA / Scripts", "Generic AI Chatbots", "TradeSentinel Copilot"]
        rows = [
            ["Workflow Creation", "Manual Scripting / Recording", "Text Advice (No Execution)", "Conversational NL -> Executable DAG"],
            ["Domain Intelligence", "None (Static Scripts)", "Generic Public Text", "Built-in ML (Risk, ETA, Customs, Routes)"],
            ["Conflict Detection", "None (Fails at runtime)", "None (Prone to errors)", "Deterministic Graph & Safety Engine"],
            ["Pre-Execution Simulation", "None (Live trial-and-error)", "None (No dataset context)", "Historical Replay (500+ Shipments)"],
            ["Human-in-the-Loop", "Basic pause scripts", "No governance gates", "Multi-Tiered Approval (> ₹10L Cargo)"],
            ["Self-Optimization", "Static rules", "None", "Auto-Approval Policy Tuner"]
        ]
        
        # Style table headers
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                r = p.runs[0]
                r.font.name = "Segoe UI"
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = CYAN
                
        # Style table rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = text
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                if p.runs:
                    r = p.runs[0]
                    r.font.name = "Segoe UI"
                    r.font.size = Pt(13)
                    r.font.bold = True if col_idx == 3 else False
                    r.font.color.rgb = GREEN if col_idx == 3 else (WHITE if col_idx == 0 else SLATE_GRAY)

# ==============================================================================
# SLIDE 9: Thank You
# ==============================================================================
slide9 = prs.slides[8]
for s in slide9.shapes:
    if s.name == 'object 2':
        tf = s.text_frame
        tf.clear()
        
        p1 = tf.paragraphs[0]
        p1.text = "Thank You!"
        p1.alignment = PP_ALIGN.CENTER
        if p1.runs:
            r = p1.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(60)
            r.font.bold = True
            r.font.color.rgb = WHITE
            
        p2 = tf.add_paragraph()
        p2.text = "TradeSentinel — Autonomous Logistics Intelligence for Global Trade"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)
        if p2.runs:
            r = p2.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = CYAN
            
        p3 = tf.add_paragraph()
        p3.text = "Team MARK42  |  PS4: AI-Powered Business Automation Copilot"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(10)
        if p3.runs:
            r = p3.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(20)
            r.font.italic = True
            r.font.color.rgb = LIGHT_BLUE
            
        p4 = tf.add_paragraph()
        p4.text = "Ready for Live Demonstration & Q&A"
        p4.alignment = PP_ALIGN.CENTER
        p4.space_before = Pt(14)
        if p4.runs:
            r = p4.runs[0]
            r.font.name = "Segoe UI"
            r.font.size = Pt(18)
            r.font.bold = True
            r.font.color.rgb = GREEN

# Save final presentation
prs.save('CIT Hackthon.pptx')
print("Successfully generated final CIT Hackthon.pptx!")
