import pptx
import os

prs = pptx.Presentation('CIT Hackthon.pptx')
os.makedirs('extracted_images', exist_ok=True)

img_count = 0
for s_idx, slide in enumerate(prs.slides):
    for sh_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.ext
            filename = f"extracted_images/slide_{s_idx+1}_shape_{sh_idx}_{shape.name.replace(' ', '_')}.{ext}"
            with open(filename, 'wb') as f:
                f.write(img.blob)
            print(f"Saved {filename} ({len(img.blob)} bytes)")
            img_count += 1

print(f"Total extracted: {img_count}")
