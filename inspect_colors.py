import pptx
from pptx.dml.color import RGBColor

prs = pptx.Presentation('CIT Hackthon.pptx')

for i, slide in enumerate(prs.slides):
    bg = slide.background
    fill = bg.fill
    print(f"Slide {i+1} background type: {fill.type if fill else 'None'}")
    
    # Check text colors and fonts in shapes
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.color and r.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                        print(f"  Slide {i+1} Text: font={r.font.name}, size={r.font.size.pt if r.font.size else None}, color=#{r.font.color.rgb}, text={r.text[:30]}")
                    elif r.font.color and r.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.THEME:
                        print(f"  Slide {i+1} Text: font={r.font.name}, theme_color={r.font.color.theme_color}, text={r.text[:30]}")
                    else:
                        print(f"  Slide {i+1} Text: font={r.font.name}, size={r.font.size.pt if r.font.size else None}, color=None, text={r.text[:30]}")
