import pptx
import json

prs = pptx.Presentation('CIT Hackthon.pptx')
slide_info = []

for idx, slide in enumerate(prs.slides):
    s_data = {
        'slide_index': idx + 1,
        'layout': slide.slide_layout.name,
        'shapes': []
    }
    for shape in slide.shapes:
        sh_data = {
            'id': shape.shape_id,
            'name': shape.name,
            'type': str(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
        }
        if shape.has_text_frame:
            paragraphs = []
            for p in shape.text_frame.paragraphs:
                p_info = {
                    'text': p.text,
                    'alignment': str(p.alignment),
                    'runs': []
                }
                for r in p.runs:
                    color_str = None
                    if r.font.color and r.font.color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB:
                        color_str = str(r.font.color.rgb)
                    p_info['runs'].append({
                        'text': r.text,
                        'font_name': r.font.name,
                        'font_size': r.font.size.pt if r.font.size else None,
                        'bold': r.font.bold,
                        'italic': r.font.italic,
                        'color': color_str
                    })
                paragraphs.append(p_info)
            sh_data['paragraphs'] = paragraphs
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            sh_data['table'] = table_data
        s_data['shapes'].append(sh_data)
    slide_info.append(s_data)

with open('template_dump.json', 'w', encoding='utf-8') as f:
    json.dump(slide_info, f, indent=2, ensure_ascii=False)

print("Saved template_dump.json successfully!")
