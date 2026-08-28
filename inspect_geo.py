import pptx

prs = pptx.Presentation('CIT Hackthon.pptx')
print(f"Slide Width: {prs.slide_width / 914400} inches, Slide Height: {prs.slide_height / 914400} inches")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    print(f"Layout: {slide.slide_layout.name}")
    print(f"Shapes count: {len(slide.shapes)}")
    for s in slide.shapes:
        print(f"  Shape: '{s.name}' (id={s.shape_id}, type={s.shape_type}) pos=({s.left/914400:.2f}, {s.top/914400:.2f}) size=({s.width/914400:.2f} x {s.height/914400:.2f} in)")
